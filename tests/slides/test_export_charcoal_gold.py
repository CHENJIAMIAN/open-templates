from pathlib import Path
import subprocess
import sys

import pytest
from pptx import Presentation

from tools.slides.sml_parser import build_asset_index, parse_presentation, resolve_asset_token
from tools.slides.pptx_writer import write_presentation


def test_export_charcoal_gold_creates_artifacts(tmp_path: Path):
    source_dir = Path("templates/slides/charcoal_gold")
    output_dir = tmp_path / "out"

    result = subprocess.run(
        [
            sys.executable,
            "tools/slides/export_charcoal_gold.py",
            "--source",
            str(source_dir),
            "--output",
            str(output_dir),
        ],
        capture_output=True,
        text=True,
    )

    assert result.returncode == 0, result.stderr
    assert (output_dir / "deck.pptx").exists()
    assert (output_dir / "deck.pdf").exists()


def test_parse_charcoal_gold_extracts_deck_structure():
    deck = parse_presentation(Path("templates/slides/charcoal_gold/template.xml"))

    assert deck.width == 960
    assert deck.height == 540
    assert len(deck.slides) == 20
    assert any(node.kind == "text" for slide in deck.slides for node in slide.nodes)
    assert any(node.kind == "image" for slide in deck.slides for node in slide.nodes)
    assert any(node.kind == "line" for slide in deck.slides for node in slide.nodes)
    assert any(node.kind == "shape" for slide in deck.slides for node in slide.nodes)


def test_asset_resolution_maps_sml_token_to_local_png():
    asset_index = build_asset_index(Path("templates/slides/charcoal_gold"))
    path = resolve_asset_token(asset_index, "FJLsbVgEEogLRtxWLW2cOe3Tn7f")

    assert path.name.endswith("FJLsbVgEEogLRtxWLW2cOe3Tn7f.png")


def test_asset_resolution_fails_for_unknown_token():
    asset_index = build_asset_index(Path("templates/slides/charcoal_gold"))

    with pytest.raises(KeyError):
        resolve_asset_token(asset_index, "missing-token")


def test_write_pptx_preserves_slide_count(tmp_path: Path):
    deck = parse_presentation(Path("templates/slides/charcoal_gold/template.xml"))
    output = tmp_path / "deck.pptx"

    write_presentation(deck, Path("templates/slides/charcoal_gold"), output)

    prs = Presentation(output)
    assert output.exists()
    assert output.stat().st_size > 0
    assert len(prs.slides) == len(deck.slides)

    for parsed_slide, generated_slide in zip(deck.slides, prs.slides):
        if parsed_slide.nodes:
            assert len(generated_slide.shapes) >= 1

    expected_images = sum(1 for slide in deck.slides for node in slide.nodes if node.kind == "image")
    expected_lines = sum(1 for slide in deck.slides for node in slide.nodes if node.kind == "line")
    expected_shapes = sum(1 for slide in deck.slides for node in slide.nodes if node.kind == "shape")
    actual_images = sum(1 for slide in prs.slides for shape in slide.shapes if shape.shape_type == 13)
    actual_freeform = sum(1 for slide in prs.slides for shape in slide.shapes if shape.shape_type in {1, 9, 17})

    assert actual_images == expected_images
    assert actual_freeform >= expected_lines + expected_shapes
