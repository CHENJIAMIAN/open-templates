from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation as PptxPresentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.util import Emu, Pt

from tools.slides.sml_parser import Presentation
from tools.slides.sml_parser import build_asset_index, resolve_asset_token


def write_presentation(deck: Presentation, template_dir: Path, output_path: Path) -> None:
    prs = PptxPresentation()
    prs.slide_width = _px(deck.width)
    prs.slide_height = _px(deck.height)
    asset_index = build_asset_index(template_dir)
    blank_layout = prs.slide_layouts[6]

    while prs.slides:
        xml_slide = prs.slides._sldIdLst[-1]
        prs.part.drop_rel(xml_slide.rId)
        del prs.slides._sldIdLst[-1]

    for slide_model in deck.slides:
        slide = prs.slides.add_slide(blank_layout)
        _apply_background(slide, slide_model.background_color)
        for node in slide_model.nodes:
            _add_node(slide, node, asset_index)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def _add_node(slide, node, asset_index: dict[str, Path]) -> None:
    if node.kind == "text":
        textbox = slide.shapes.add_textbox(_px(node.x), _px(node.y), _px(node.width), _px(node.height))
        textbox.text = node.text or ""
        if node.fill_color:
            textbox.fill.solid()
            textbox.fill.fore_color.rgb = _rgb(node.fill_color)
        else:
            textbox.fill.background()
        for paragraph in textbox.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(18)
    elif node.kind == "image" and node.src:
        image_path = resolve_asset_token(asset_index, node.src)
        slide.shapes.add_picture(str(image_path), _px(node.x), _px(node.y), _px(node.width), _px(node.height))
    elif node.kind == "line":
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            _px(node.start_x),
            _px(node.start_y),
            _px(node.end_x),
            _px(node.end_y),
        )
        if node.border_color:
            line.line.color.rgb = _rgb(node.border_color)
        if node.border_width is not None:
            line.line.width = _px(node.border_width)
    elif node.kind == "shape":
        shape = slide.shapes.add_shape(
            _shape_type(node.shape_type),
            _px(node.x),
            _px(node.y),
            _px(node.width),
            _px(node.height),
        )
        if node.fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(node.fill_color)
        else:
            shape.fill.background()
        shape.line.fill.background()


def _apply_background(slide, color: str | None) -> None:
    if not color:
        return
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _shape_type(shape_type: str | None) -> MSO_AUTO_SHAPE_TYPE:
    mapping = {
        "rect": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        "round-rect": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        "diamond": MSO_AUTO_SHAPE_TYPE.DIAMOND,
        "ellipse": MSO_AUTO_SHAPE_TYPE.OVAL,
    }
    if shape_type not in mapping:
        raise ValueError(f"unsupported shape type: {shape_type}")
    return mapping[shape_type]


def _px(value: float | int | None) -> Emu:
    if value is None:
        return Emu(0)
    return Emu(round(float(value) * 914400 / 96))


def _rgb(color: str) -> RGBColor:
    match = re.fullmatch(r"rgba\(\s*(\d+),\s*(\d+),\s*(\d+),\s*([0-9.]+)\s*\)", color)
    if not match:
        raise ValueError(f"unsupported color format: {color}")
    red, green, blue = (int(match.group(index)) for index in range(1, 4))
    return RGBColor(red, green, blue)
